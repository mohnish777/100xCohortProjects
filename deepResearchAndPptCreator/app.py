import os
from typing import List, Optional
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from groq import Groq
from serpapi import GoogleSearch
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
import gradio as gr
import asyncio

# Load environment variables
load_dotenv()

# Initialize FastAPI app
app = FastAPI(title="Research and Presentation Generator")

# Initialize Groq client
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

class ResearchRequest(BaseModel):
    topic: str
    num_slides: int = 5

class ResearchResponse(BaseModel):
    research_content: str
    slides_created: bool
    file_path: Optional[str] = None

def perform_research(topic: str) -> str:
    """Perform research using SerpAPI and Groq."""
    try:
        # Search using SerpAPI
        search = GoogleSearch({
            "q": topic,
            "api_key": os.getenv("SERPAPI_API_KEY"),
            "num": 10
        })
        results = search.get_dict()
        
        # Extract relevant information
        search_results = []
        if "organic_results" in results:
            for result in results["organic_results"]:
                search_results.append(f"Title: {result.get('title', '')}\nSnippet: {result.get('snippet', '')}\n")
        
        # Use Groq to analyze and synthesize the research
        prompt = f"""Based on the following research about {topic}, create a comprehensive analysis:
        
        Research findings:
        {' '.join(search_results)}
        
        Please provide a detailed analysis that can be used for a presentation."""
        
        chat_completion = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="meta-llama/llama-4-maverick-17b-128e-instruct",
        )
        
        return chat_completion.choices[0].message.content
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Research failed: {str(e)}")

def create_presentation(content: str, num_slides: int, topic: str) -> str:
    """Create a PowerPoint presentation from the research content."""
    try:
        prs = Presentation()
        
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = topic
        subtitle.text = "Research Presentation"
        
        # Split content into sections for slides
        sections = content.split('\n\n')
        for i, section in enumerate(sections[:num_slides]):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"Section {i+1}"
            content.text = section
        
        # Save the presentation
        output_path = f"presentations/{topic.replace(' ', '_')}_presentation.pptx"
        os.makedirs("presentations", exist_ok=True)
        prs.save(output_path)
        return output_path
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Presentation creation failed: {str(e)}")

#@app.post("/research", response_model=ResearchResponse)
async def research_topic(request: ResearchRequest):
    """Endpoint to perform research and create presentation."""
    try:
        # Perform research
        research_content = perform_research(request.topic)
        
        # Create presentation
        file_path = create_presentation(research_content, request.num_slides, request.topic)
        
        return ResearchResponse(
            research_content=research_content,
            slides_created=True,
            file_path=file_path
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Gradio Interface
async def gradio_interface(topic: str, num_slides: int):
    """Gradio interface for the research and presentation generator."""
    try:
        request = ResearchRequest(topic=topic, num_slides=num_slides)
        response = await research_topic(request)
        return response.research_content, response.file_path
    except Exception as e:
        return str(e), None

# Create Gradio interface
iface = gr.Interface(
    fn=gradio_interface,
    inputs=[
        gr.Textbox(label="Research Topic"),
        gr.Slider(minimum=1, maximum=10, value=5, step=1, label="Number of Slides")
    ],
    outputs=[
        gr.Textbox(label="Research Content"),
        gr.File(label="Download Presentation")
    ],
    title="Research and Presentation Generator",
    description="Enter a topic to research and generate a presentation."
)

# Mount Gradio app to FastAPI
app = gr.mount_gradio_app(app, iface, path="/")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
